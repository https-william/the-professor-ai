import os
import tempfile
import base64
import io
import fitz  # PyMuPDF
from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from werkzeug.utils import secure_filename

app = Flask(__name__)

def extract_images_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    extracted_images = []

    for page_index in range(len(doc)):
        page = doc[page_index]
        image_list = page.get_images(full=True)
        
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            # Convert to base64 for frontend
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            extracted_images.append({
                "page": page_index + 1,
                "index": img_index,
                "data": base64_image,
                "ext": base_image["ext"]
            })
            
    doc.close()
    return extracted_images

def extract_images_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    extracted_images = []

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                base64_image = base64.b64encode(image_bytes).decode('utf-8')
                extracted_images.append({
                    "page": slide_index + 1,
                    "index": shape_index,
                    "data": base64_image,
                    "ext": shape.image.ext
                })
    return extracted_images

@app.route('/api/extract_images', methods=['POST'])
@app.route('/', methods=['POST'])
def handle_extract():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    filename = secure_filename(file.filename)
    file_ext = os.path.splitext(filename)[1].lower()
    
    with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as temp:
        file.save(temp.name)
        temp_path = temp.name
    
    try:
        images = []
        if file_ext == '.pdf':
            images = extract_images_from_pdf(temp_path)
        elif file_ext == '.pptx':
            images = extract_images_from_pptx(temp_path)
        else:
            return jsonify({"error": "Unsupported file type for image extraction"}), 400
            
        return jsonify({
            "success": True,
            "images": images,
            "count": len(images)
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500
        
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

def handler(request):
    return app(request)

if __name__ == "__main__":
    app.run(port=5001)
